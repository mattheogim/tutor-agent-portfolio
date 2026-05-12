"""Convert all .pptx files in codenotes/ to _slides.md markdown files."""

import re
from datetime import date
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

CODENOTES_DIR = Path(__file__).resolve().parent.parent / "codenotes"

def looks_like_code(text: str) -> bool:
    """Heuristic: detect C++ code content."""
    indicators = [
        "int ", "void ", "double ", "char ", "bool ",
        "cout", "cin", "endl", "#include", "return ",
        "class ", "struct ", "template", "nullptr",
        "{", "}", "//", "/*", "->", "::",
    ]
    score = sum(1 for ind in indicators if ind in text)
    return score >= 3

def has_monospace_font(run) -> bool:
    """Check if a run uses a monospace font."""
    mono_fonts = {"courier", "consolas", "monaco", "menlo", "monospace", "source code"}
    if run.font and run.font.name:
        return any(m in run.font.name.lower() for m in mono_fonts)
    return False

def extract_table(table) -> str:
    """Convert a pptx table to markdown table."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    lines = []
    # Header
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    for row in rows[1:]:
        # Pad if needed
        while len(row) < len(rows[0]):
            row.append("")
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)

def extract_text_frame(text_frame) -> str:
    """Extract text from a text frame, preserving structure."""
    blocks = []
    code_lines = []
    in_code = False

    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            if in_code:
                code_lines.append("")
            continue

        # Check if this paragraph looks like code
        is_code = any(has_monospace_font(run) for run in para.runs if run.text.strip())

        if is_code:
            if not in_code:
                in_code = True
                code_lines = []
            code_lines.append(para.text.rstrip())
        else:
            if in_code:
                blocks.append("```cpp\n" + "\n".join(code_lines) + "\n```")
                in_code = False
                code_lines = []

            level = para.level if para.level else 0
            indent = "  " * level
            if level > 0:
                blocks.append(f"{indent}- {text}")
            else:
                blocks.append(text)

    if in_code and code_lines:
        blocks.append("```cpp\n" + "\n".join(code_lines) + "\n```")

    return "\n".join(blocks)

def convert_slide(slide, slide_num: int) -> str:
    """Convert a single slide to markdown."""
    title = ""
    content_parts = []

    for shape in slide.shapes:
        if shape.has_table:
            content_parts.append(extract_table(shape.table))
        elif shape.has_text_frame:
            is_title = False
            try:
                if shape.placeholder_format and shape.placeholder_format.idx in (0, 1):
                    is_title = True
            except ValueError:
                pass
            if is_title and not title:
                title = shape.text_frame.text.strip()
            else:
                extracted = extract_text_frame(shape.text_frame)
                if extracted:
                    content_parts.append(extracted)
        elif hasattr(shape, "image"):
            alt = shape.name or "slide visual"
            content_parts.append(f"[Image: {alt}]")

    # If no title found from placeholder, check if first shape is title-like
    if not title:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                candidate = shape.text_frame.text.strip()
                if len(candidate) < 80 and "\n" not in candidate:
                    title = candidate
                    # Remove from content if it was also added there
                    content_parts = [p for p in content_parts if p.strip() != candidate]
                    break

    header = f"## Slide {slide_num}" + (f": {title}" if title else "")
    body = "\n\n".join(content_parts) if content_parts else "(empty slide)"

    # Post-process: if a whole block looks like code but wasn't detected
    if looks_like_code(body) and "```" not in body:
        lines = body.split("\n")
        if all(not l.startswith("- ") and not l.startswith("| ") for l in lines if l.strip()):
            body = f"```cpp\n{body}\n```"

    return f"{header}\n\n{body}"

def convert_pptx(pptx_path: Path) -> str:
    """Convert a .pptx file to markdown string."""
    prs = Presentation(str(pptx_path))
    slides_md = []

    for i, slide in enumerate(prs.slides, 1):
        slides_md.append(convert_slide(slide, i))

    header = (
        f"# {pptx_path.stem} — Lecture Slides\n\n"
        f"> Source: `{pptx_path.name}` | Converted: {date.today()} | DO NOT EDIT — regenerate from source\n"
    )

    return header + "\n\n---\n\n".join([""] + slides_md) + "\n"

def main():
    pptx_files = sorted(CODENOTES_DIR.rglob("*.pptx"))
    print(f"Found {len(pptx_files)} .pptx files\n")

    for pptx_path in pptx_files:
        out_path = pptx_path.with_name(pptx_path.stem + "_slides.md")
        try:
            md = convert_pptx(pptx_path)
            out_path.write_text(md, encoding="utf-8")
            slides_count = md.count("\n## Slide ")
            print(f"  ✅ {pptx_path.relative_to(CODENOTES_DIR)} → {out_path.name} ({slides_count} slides)")
        except Exception as e:
            print(f"  ❌ {pptx_path.relative_to(CODENOTES_DIR)} — {e}")

    print(f"\nDone. Original .pptx files preserved.")

if __name__ == "__main__":
    main()
